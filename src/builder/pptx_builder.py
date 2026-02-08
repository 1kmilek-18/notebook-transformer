"""python-pptxを使用してPowerPointファイルを構築するモジュール."""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from src.models import (
    ElementType,
    ImageBlock,
    PresentationData,
    SlideData,
    TextBlock,
    TextSpan,
)
from src.utils.coordinate import pt_to_emu

logger = logging.getLogger(__name__)

# テキスト揃えのマッピング
ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


class PPTXBuilder:
    """PresentationDataからPowerPointファイルを構築するクラス.

    抽出・解析済みのスライドデータを元に、python-pptxを使って
    座標を維持しながらテキストボックスと画像を配置する。
    """

    def __init__(self, template_path: Optional[str | Path] = None) -> None:
        """PPTXBuilderを初期化する.

        Args:
            template_path: PowerPointテンプレート(.potx/.pptx)のパス。
                          Noneの場合は空のプレゼンテーションを作成。
        """
        self.template_path = Path(template_path) if template_path else None
        self._prs: Optional[Presentation] = None

    def build(self, data: PresentationData) -> Presentation:
        """PresentationDataからPowerPointプレゼンテーションを構築する.

        スライドサイズ・各スライドのテキストボックス・画像を座標維持で配置する。
        テンプレートが指定されていればそのマスターを利用する。

        Args:
            data: 抽出・解析済みプレゼンテーションデータ

        Returns:
            python-pptx の Presentation オブジェクト（save() で保存するまでメモリ上のみ）
        """
        logger.info("PowerPoint構築を開始: %d スライド", len(data.slides))

        # プレゼンテーション初期化
        if self.template_path and self.template_path.exists():
            logger.info("テンプレートを使用: %s", self.template_path)
            self._prs = Presentation(str(self.template_path))
        else:
            self._prs = Presentation()

        # スライドサイズ設定（PDF座標に合わせる）
        self._prs.slide_width = Emu(pt_to_emu(data.slide_width))
        self._prs.slide_height = Emu(pt_to_emu(data.slide_height))

        # 各スライドを構築
        for slide_data in data.slides:
            self._build_slide(slide_data)

        logger.info("PowerPoint構築完了")
        return self._prs

    def save(self, output_path: str | Path) -> Path:
        """構築したプレゼンテーションを保存する.

        出力先の親ディレクトリが存在しない場合は自動作成する。
        既存ファイルは上書きする。

        Args:
            output_path: 保存先ファイルパス（.pptx）

        Returns:
            保存されたファイルの Path

        Raises:
            RuntimeError: build() が呼ばれていない場合
            OSError: ディレクトリ作成またはファイル書き込みに失敗した場合
        """
        if self._prs is None:
            raise RuntimeError("プレゼンテーションが構築されていません。build()を先に呼び出してください。")

        path = Path(output_path)
        parent = path.parent
        existed = parent.exists()
        try:
            parent.mkdir(parents=True, exist_ok=True)
        except OSError:
            logger.error("出力ディレクトリの作成に失敗: %s", parent, exc_info=True)
            raise
        if not existed:
            logger.info("出力ディレクトリを作成: %s", parent)
        self._prs.save(str(path))
        logger.info("PowerPointを保存: %s", path)
        return path

    def _build_slide(self, slide_data: SlideData) -> None:
        """1スライド分を構築する.

        Args:
            slide_data: 1スライド分のデータ
        """
        assert self._prs is not None

        # 空白レイアウトでスライド追加
        blank_layout = self._prs.slide_layouts[6]  # 空白レイアウト
        slide = self._prs.slides.add_slide(blank_layout)

        # テキストブロックを配置
        for text_block in slide_data.text_blocks:
            self._add_text_box(slide, text_block)

        # 画像を配置
        for image_block in slide_data.image_blocks:
            self._add_image(slide, image_block)

        logger.debug(
            "スライド %d: テキスト %d個, 画像 %d個",
            slide_data.page_number,
            len(slide_data.text_blocks),
            len(slide_data.image_blocks),
        )

    def _add_text_box(self, slide: object, block: TextBlock) -> None:
        """スライドにテキストボックスを追加する.

        Args:
            slide: python-pptxのSlideオブジェクト
            block: テキストブロックデータ
        """
        left = Emu(pt_to_emu(block.bbox.x0))
        top = Emu(pt_to_emu(block.bbox.y0))
        width = Emu(pt_to_emu(block.bbox.width))
        height = Emu(pt_to_emu(block.bbox.height))

        # テキストボックスを追加
        txBox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
        tf = txBox.text_frame
        tf.word_wrap = True

        # 最初のスパンは既存の段落に追加
        first_span = True
        current_paragraph = tf.paragraphs[0]

        # テキスト揃えを設定
        alignment = ALIGNMENT_MAP.get(block.alignment, PP_ALIGN.LEFT)
        current_paragraph.alignment = alignment

        for span in block.spans:
            if not span.text:
                continue

            if first_span:
                run = current_paragraph.add_run()
                first_span = False
            else:
                # 改行を含む場合は新しい段落を追加
                if "\n" in span.text:
                    lines = span.text.split("\n")
                    for i, line in enumerate(lines):
                        if i > 0:
                            current_paragraph = tf.add_paragraph()
                            current_paragraph.alignment = alignment
                        if line.strip():
                            run = current_paragraph.add_run()
                            run.text = line
                            self._apply_font(run, span)
                    continue
                else:
                    run = current_paragraph.add_run()

            run.text = span.text
            self._apply_font(run, span)

    def _apply_font(self, run: object, span: TextSpan) -> None:
        """Runオブジェクトにフォント設定を適用する.

        Args:
            run: python-pptxのRunオブジェクト
            span: フォント情報を含むTextSpan
        """
        font = run.font  # type: ignore[attr-defined]
        font.size = Pt(span.font.size)
        font.bold = span.font.bold
        font.italic = span.font.italic

        # フォント名設定
        if span.font.name:
            font.name = span.font.name

        # 色設定
        try:
            color_hex = span.font.color.lstrip("#")
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            font.color.rgb = RGBColor(r, g, b)
        except (ValueError, IndexError) as e:
            logger.warning("色の変換に失敗: %s (%s)", span.font.color, e)

    def _add_image(self, slide: object, image_block: ImageBlock) -> None:
        """スライドに画像を追加する.

        Args:
            slide: python-pptxのSlideオブジェクト
            image_block: 画像ブロックデータ
        """
        left = Emu(pt_to_emu(image_block.bbox.x0))
        top = Emu(pt_to_emu(image_block.bbox.y0))
        width = Emu(pt_to_emu(image_block.bbox.width))
        height = Emu(pt_to_emu(image_block.bbox.height))

        if image_block.source_path:
            # ファイルパスから画像を追加
            slide.shapes.add_picture(  # type: ignore[attr-defined]
                image_block.source_path, left, top, width, height
            )
        elif image_block.image_data:
            # バイトデータから画像を追加
            image_stream = io.BytesIO(image_block.image_data)
            slide.shapes.add_picture(  # type: ignore[attr-defined]
                image_stream, left, top, width, height
            )
        else:
            logger.warning("画像データが空です")
